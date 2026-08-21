"""Auxiliary resources for the Responses surface."""

from __future__ import annotations

import base64
import hashlib
import io
import json
import uuid
from datetime import UTC, datetime
from pathlib import Path

from fastapi import APIRouter, Depends, File, Form, Request, UploadFile
from sqlalchemy import func, select
from sqlalchemy.ext.asyncio import AsyncSession

from data_agent.adapters.opentrace.learning import OpenTraceLearningRepository
from data_agent.adapters.opentrace.repository import OpenTraceRunRepository
from data_agent.contracts import DataScope
from gateway.api_gateway.routers.auth import get_current_user
from gateway.api_gateway.tenant_middleware import build_tenant_metadata
from infra.config.settings import settings
from infra.errors import AppException, ErrorCodes
from infra.observability.metrics import (
    RESPONSE_APPROVAL_DECISIONS_TOTAL,
    RESPONSE_APPROVAL_RESOLUTION_DURATION,
)
from infra.responses.repository import TERMINAL_STATUSES, append_event
from infra.security.identity import is_enterprise_admin
from infra.storage.data_agent_models import DataAgentFeedback, DataAgentRunRecord
from infra.storage.database import db_session_dependency as get_db
from infra.storage.models import (
    Attachment,
    ChatSession,
    Feedback,
    ResponseApproval,
    ResponseEvent,
    ResponseItem,
    ResponseOutbox,
    ResponseRecord,
    ResponseToolExecution,
    User,
)
from infra.storage.object_store import attachment_object_key, get_object_store
from services.production_intelligence.audit import mask_sensitive

router = APIRouter()

_TEXT_EXTENSIONS = {
    ".csv",
    ".css",
    ".html",
    ".js",
    ".json",
    ".log",
    ".md",
    ".py",
    ".txt",
    ".ts",
    ".xml",
    ".yaml",
    ".yml",
}
_DOCUMENT_EXTENSIONS = {".docx", ".pdf", ".pptx", ".xlsx"}
_IMAGE_EXTENSIONS = {".gif", ".jpeg", ".jpg", ".png", ".webp"}
_AUDIO_EXTENSIONS = {".aac", ".flac", ".m4a", ".mp3", ".ogg", ".wav"}
_VIDEO_EXTENSIONS = {".avi", ".flv", ".mkv", ".mov", ".mp4", ".webm", ".wmv"}
_MEDIA_MIME_BY_EXTENSION = {
    ".aac": "audio/aac",
    ".flac": "audio/flac",
    ".m4a": "audio/mp4",
    ".mp3": "audio/mpeg",
    ".ogg": "audio/ogg",
    ".wav": "audio/wav",
    ".avi": "video/x-msvideo",
    ".flv": "video/x-flv",
    ".mkv": "video/x-matroska",
    ".mov": "video/quicktime",
    ".mp4": "video/mp4",
    ".webm": "video/webm",
    ".wmv": "video/x-ms-wmv",
}


async def _extract_attachment_text(raw: bytes, filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        try:
            lines: list[str] = []
            for sheet in workbook.worksheets:
                lines.append(f"[工作表: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        lines.append("\t".join(values))
            return "\n".join(lines)
        finally:
            workbook.close()
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(raw))
        lines = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"[幻灯片 {index}]")
            lines.extend(
                str(shape.text).strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and str(shape.text).strip()
            )
        return "\n".join(lines)
    from gateway.api_gateway.routers.documents import _extract_text

    return await _extract_text(raw, filename)


def _attachment_out(row: Attachment) -> dict:
    return {
        "id": row.id,
        "attachment_id": row.id,
        "filename": row.filename,
        "file_size": int(row.file_size or 0),
        "mime_type": row.mime_type,
        "file_extension": row.file_extension,
        "content_summary": row.content_summary,
        "content_hash": row.content_hash,
        "status": row.status,
        "message_id": row.message_id,
        "scope": row.scope,
        "ingest_status": row.ingest_status,
        "promoted_document_id": row.promoted_document_id,
        "media_kind": row.media_kind,
        "storage_backend": row.storage_backend,
        "created_at": row.created_at.isoformat() if row.created_at else None,
    }


async def _owned_session(
    db: AsyncSession, *, session_id: str, user: User, request: Request
) -> ChatSession:
    tenant_id, workspace_id = _scope(request, user)
    row = await db.scalar(
        select(ChatSession).where(
            ChatSession.id == session_id,
            ChatSession.user_id == user.id,
            ChatSession.tenant_id == tenant_id,
            ChatSession.workspace_id == workspace_id,
        )
    )
    if row is None:
        raise AppException(
            ErrorCodes.RESOURCE_NOT_FOUND.code, message="Conversation 不存在或无权限"
        )
    return row


@router.post("/files")
async def upload_response_file(
    request: Request,
    file: UploadFile = File(...),
    session_id: str = Form(...),
    message_id: str | None = Form(default=None),
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> dict:
    session = await _owned_session(db, session_id=session_id, user=current_user, request=request)
    raw = await file.read()
    max_bytes = max(1, int(settings.attachment_max_size_mb)) * 1024 * 1024
    if not raw:
        raise AppException(ErrorCodes.PARAM_INVALID.code, message="附件不能为空")
    if len(raw) > max_bytes:
        raise AppException(
            ErrorCodes.PARAM_INVALID.code,
            message=f"附件不能超过 {settings.attachment_max_size_mb}MB",
        )
    filename = Path(file.filename or "attachment").name[:512]
    suffix = Path(filename).suffix.lower()
    supported = (
        _TEXT_EXTENSIONS
        | _DOCUMENT_EXTENSIONS
        | _IMAGE_EXTENSIONS
        | _AUDIO_EXTENSIONS
        | _VIDEO_EXTENSIONS
    )
    if suffix not in supported:
        raise AppException(
            ErrorCodes.PARAM_INVALID.code,
            message="暂不支持该附件格式",
        )
    content_hash = hashlib.sha256(raw).hexdigest()
    duplicate = await db.scalar(
        select(Attachment).where(
            Attachment.session_id == session_id,
            Attachment.user_id == current_user.id,
            Attachment.content_hash == content_hash,
            Attachment.status == "active",
        )
    )

    mime_type = str(file.content_type or "application/octet-stream")[:255]
    is_image = mime_type.startswith("image/")
    is_audio = mime_type.startswith("audio/") or suffix in _AUDIO_EXTENSIONS
    is_video = mime_type.startswith("video/") or suffix in _VIDEO_EXTENSIONS
    if suffix in _IMAGE_EXTENSIONS:
        is_image = True
        mime_type = {
            ".gif": "image/gif",
            ".jpeg": "image/jpeg",
            ".jpg": "image/jpeg",
            ".png": "image/png",
            ".webp": "image/webp",
        }[suffix]
    elif suffix in _MEDIA_MIME_BY_EXTENSION:
        mime_type = _MEDIA_MIME_BY_EXTENSION[suffix]
    if is_audio or is_video:
        media_limit = max(1, int(settings.multimodal_inline_max_mb)) * 1024 * 1024
        if len(raw) > media_limit:
            raise AppException(
                ErrorCodes.PARAM_INVALID.code,
                message=(
                    f"音视频附件需小于 {settings.multimodal_inline_max_mb}MB，"
                    "以便作为原生多模态内容安全传入模型"
                ),
            )
    content_text = ""
    if not (is_image or is_audio or is_video):
        try:
            content_text = (await _extract_attachment_text(raw, filename)).strip()
        except Exception as exc:
            raise AppException(
                ErrorCodes.PARAM_INVALID.code,
                message="附件无法解析或文件已损坏",
            ) from exc
    media_kind = "image" if is_image else "audio" if is_audio else "video" if is_video else None
    media_label = {"image": "图片", "audio": "音频", "video": "视频"}.get(media_kind or "", "文件")
    summary = content_text[:512] if content_text else f"{media_label}附件：{filename}"
    object_store = get_object_store()
    object_key = None
    object_etag = None
    storage_backend = "database"
    if object_store is not None:
        if duplicate and duplicate.object_key:
            object_key = duplicate.object_key
            object_etag = duplicate.object_etag
            storage_backend = duplicate.storage_backend
        else:
            object_key = attachment_object_key(
                tenant_id=session.tenant_id,
                workspace_id=session.workspace_id,
                content_hash=content_hash,
            )
            object_ref = await object_store.put(object_key, raw, mime_type)
            object_etag = object_ref.etag
            storage_backend = object_ref.backend
    row = Attachment(
        id=str(uuid.uuid4()),
        session_id=session_id,
        user_id=current_user.id,
        tenant_id=session.tenant_id,
        workspace_id=session.workspace_id,
        filename=filename,
        file_size=len(raw),
        mime_type=mime_type,
        file_extension=Path(filename).suffix.lower().lstrip(".")[:20] or None,
        content_hash=content_hash,
        duplicate_of=duplicate.id if duplicate else None,
        content_text=content_text[: max(1, int(settings.attachment_max_chars)) * 8] or None,
        content_summary=summary,
        image_base64=(
            base64.b64encode(raw).decode("ascii") if is_image and object_store is None else None
        ),
        image_mime=mime_type if is_image else None,
        media_base64=(
            base64.b64encode(raw).decode("ascii")
            if (is_audio or is_video) and object_store is None
            else None
        ),
        media_mime=mime_type if is_audio or is_video else None,
        media_kind=media_kind,
        storage_backend=storage_backend,
        object_key=object_key,
        object_etag=object_etag,
        message_id=message_id,
        status="active",
        scope="session",
        ingest_status="ready",
    )
    db.add(row)
    await db.commit()
    await db.refresh(row)
    result = _attachment_out(row)
    result["is_duplicate"] = duplicate is not None
    return result


@router.get("/files/{session_id}")
async def list_response_files(
    session_id: str,
    request: Request,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> dict:
    await _owned_session(db, session_id=session_id, user=current_user, request=request)
    rows = (
        (
            await db.execute(
                select(Attachment)
                .where(
                    Attachment.session_id == session_id,
                    Attachment.user_id == current_user.id,
                    Attachment.status == "active",
                )
                .order_by(Attachment.created_at)
            )
        )
        .scalars()
        .all()
    )
    return {
        "session_id": session_id,
        "attachments": [_attachment_out(row) for row in rows],
        "total": len(rows),
    }


@router.delete("/files/{attachment_id}")
async def delete_response_file(
    attachment_id: str,
    request: Request,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> dict:
    row = await db.scalar(
        select(Attachment).where(
            Attachment.id == attachment_id, Attachment.user_id == current_user.id
        )
    )
    if row is None:
        raise AppException(ErrorCodes.RESOURCE_NOT_FOUND.code, message="附件不存在")
    await _owned_session(db, session_id=row.session_id, user=current_user, request=request)
    row.status = "deleted"
    row.image_base64 = None
    row.media_base64 = None
    row.content_text = None
    await db.commit()
    return {"attachment_id": attachment_id, "status": "deleted"}


def _scope(request: Request, user: User) -> tuple[str, str]:
    metadata = build_tenant_metadata(request, user_id=user.id)
    return str(metadata.get("tenant_id") or "default"), str(
        metadata.get("workspace_id") or "default"
    )


async def _owned_response(
    db: AsyncSession,
    *,
    response_id: str,
    user: User,
    request: Request,
    for_update: bool = False,
) -> ResponseRecord | None:
    tenant_id, workspace_id = _scope(request, user)
    statement = select(ResponseRecord).where(
        ResponseRecord.id == response_id,
        ResponseRecord.user_id == user.id,
        ResponseRecord.tenant_id == tenant_id,
        ResponseRecord.workspace_id == workspace_id,
    )
    if for_update:
        statement = statement.with_for_update().execution_options(populate_existing=True)
    return await db.scalar(statement)


async def _workspace_response_for_approver(
    db: AsyncSession,
    *,
    response_id: str,
    user: User,
    request: Request,
) -> ResponseRecord | None:
    """四眼审批只允许同 Scope 的 SRE/Admin 获取待审批 Response 行。"""

    normalized_role = str(getattr(user, "role", "") or "").strip().lower()
    if not (is_enterprise_admin(user) or normalized_role == "sre"):
        return None
    tenant_id, workspace_id = _scope(request, user)
    return await db.scalar(
        select(ResponseRecord)
        .where(
            ResponseRecord.id == response_id,
            ResponseRecord.tenant_id == tenant_id,
            ResponseRecord.workspace_id == workspace_id,
        )
        .with_for_update()
        .execution_options(populate_existing=True)
    )


def _prepare_response_for_approval_resume(response: ResponseRecord) -> None:
    """重新排队审批结果，并保证 Worker 至少还能领取一次。"""

    response.status = "queued"
    response.completed_at = None
    response.lease_owner = None
    response.lease_expires_at = None
    response.heartbeat_at = None
    response.max_attempts = max(
        int(response.max_attempts or 0),
        int(response.attempt_count or 0) + 1,
    )


async def _ensure_approval_resume_outbox(
    db: AsyncSession,
    *,
    response_id: str,
    approval_id: str,
) -> ResponseOutbox:
    """复用审批唤醒记录，避免网络重试触发唯一键冲突。"""

    event_key = f"response.execute:{response_id}:approval-{approval_id}"
    existing = await db.scalar(
        select(ResponseOutbox)
        .where(ResponseOutbox.event_key == event_key)
        .with_for_update()
        .execution_options(populate_existing=True)
    )
    if existing is not None:
        if existing.status != "pending":
            existing.status = "pending"
            existing.available_at = datetime.now(UTC)
            existing.published_at = None
            existing.last_error = None
        return existing
    row = ResponseOutbox(
        id=f"outbox_{uuid.uuid4().hex}",
        event_key=event_key,
        aggregate_id=response_id,
        aggregate_type="response",
        event_type="response.execute",
        payload={"response_id": response_id},
    )
    db.add(row)
    return row


def _raise_approval_state_conflict(
    *,
    response: ResponseRecord,
    approval: ResponseApproval,
    requested_status: str,
    message: str,
) -> None:
    raise AppException(
        ErrorCodes.RESOURCE_EXISTS.code,
        message=message,
        details={
            "response_status": response.status,
            "approval_status": approval.status,
            "requested_status": requested_status,
        },
    )


async def _resolve_response_tool_approval(
    *,
    response_id: str,
    request: Request,
    payload: dict,
    current_user: User,
    db: AsyncSession,
    call_id: str | None = None,
    approval_id: str | None = None,
) -> dict:
    response = await _owned_response(
        db,
        response_id=response_id,
        user=current_user,
        request=request,
        for_update=True,
    )
    secondary_approver = False
    if response is None:
        response = await _workspace_response_for_approver(
            db,
            response_id=response_id,
            user=current_user,
            request=request,
        )
        secondary_approver = response is not None
    if response is None:
        raise AppException(ErrorCodes.RESOURCE_NOT_FOUND.code, message="Response 不存在或无权限")

    approval_statement = select(ResponseApproval).where(ResponseApproval.response_id == response_id)
    if approval_id is not None:
        approval_statement = approval_statement.where(ResponseApproval.id == approval_id)
    else:
        approval_statement = approval_statement.where(ResponseApproval.call_id == call_id)
    approval = await db.scalar(
        approval_statement.with_for_update().execution_options(populate_existing=True)
    )
    if approval is None:
        message = "审批记录不存在" if approval_id is not None else "工具审批记录不存在"
        raise AppException(ErrorCodes.RESOURCE_NOT_FOUND.code, message=message)
    required_approvals = max(1, min(2, int(getattr(approval, "required_approvals", 1) or 1)))
    if secondary_approver and required_approvals < 2:
        raise AppException(ErrorCodes.RESOURCE_NOT_FOUND.code, message="审批记录不存在")
    decisions = [
        dict(item)
        for item in (getattr(approval, "approval_decisions", None) or [])
        if isinstance(item, dict) and str(item.get("user_id") or "")
    ]

    approved = bool(payload.get("approved", False))
    requested_status = "approved" if approved else "rejected"
    if response.status in TERMINAL_STATUSES:
        _raise_approval_state_conflict(
            response=response,
            approval=approval,
            requested_status=requested_status,
            message="Response 已结束，不能再处理旧审批",
        )

    if approval.status in {"approved", "rejected"}:
        if approval.status != requested_status:
            _raise_approval_state_conflict(
                response=response,
                approval=approval,
                requested_status=requested_status,
                message="审批已经完成，不能更改决定",
            )
        latest_sequence = await db.scalar(
            select(func.max(ResponseEvent.sequence_number)).where(
                ResponseEvent.response_id == response_id
            )
        )
        if approval.status == "approved" and response.status == "requires_action":
            _prepare_response_for_approval_resume(response)
            await _ensure_approval_resume_outbox(
                db,
                response_id=response_id,
                approval_id=approval.id,
            )
            await db.commit()
        return {
            "approved": approval.status == "approved",
            "status": approval.status,
            "call_id": approval.call_id,
            "approval_id": approval.id,
            "required_approvals": required_approvals,
            "received_approvals": len(
                {str(item["user_id"]) for item in decisions if item.get("approved") is True}
            ),
            "starting_after": int(latest_sequence if latest_sequence is not None else -1),
        }

    if approval.status not in {"pending", "pending_secondary"}:
        _raise_approval_state_conflict(
            response=response,
            approval=approval,
            requested_status=requested_status,
            message="审批状态不允许处理",
        )
    if response.status != "requires_action":
        _raise_approval_state_conflict(
            response=response,
            approval=approval,
            requested_status=requested_status,
            message="Response 当前不在等待审批状态",
        )

    existing_decision = next(
        (item for item in decisions if str(item.get("user_id")) == current_user.id), None
    )
    if existing_decision is not None:
        if bool(existing_decision.get("approved")) != approved:
            _raise_approval_state_conflict(
                response=response,
                approval=approval,
                requested_status=requested_status,
                message="当前审批人已经提交决定，不能更改",
            )
        latest_sequence = await db.scalar(
            select(func.max(ResponseEvent.sequence_number)).where(
                ResponseEvent.response_id == response_id
            )
        )
        return {
            "approved": approval.status == "approved",
            "status": approval.status,
            "call_id": approval.call_id,
            "approval_id": approval.id,
            "required_approvals": required_approvals,
            "received_approvals": len(
                {str(item["user_id"]) for item in decisions if item.get("approved") is True}
            ),
            "current_user_decision": "approved" if approved else "rejected",
            "starting_after": int(latest_sequence if latest_sequence is not None else -1),
        }

    now = datetime.now(UTC)
    decisions.append(
        {
            "user_id": current_user.id,
            "approved": approved,
            "reason": None if approved else str(payload.get("reason") or "user rejected tool call"),
            "decided_at": now.isoformat(),
        }
    )
    approval.approval_decisions = decisions
    received_approvals = len(
        {str(item["user_id"]) for item in decisions if item.get("approved") is True}
    )
    if approved and received_approvals < required_approvals:
        approval.status = "pending_secondary"
        progress_event = await append_event(
            db,
            response_id=response_id,
            event_type="opentrace.approval.progress",
            payload={
                "approval_id": approval.id,
                "call_id": approval.call_id,
                "status": approval.status,
                "required_approvals": required_approvals,
                "received_approvals": received_approvals,
            },
        )
        await db.commit()
        RESPONSE_APPROVAL_DECISIONS_TOTAL.labels(
            outcome="pending_secondary",
            required_approvals=str(required_approvals),
        ).inc()
        return {
            "approved": False,
            "status": approval.status,
            "call_id": approval.call_id,
            "approval_id": approval.id,
            "required_approvals": required_approvals,
            "received_approvals": received_approvals,
            "current_user_decision": "approved",
            "starting_after": progress_event.sequence_number,
        }

    approval.status = requested_status
    approval.reason = None if approved else str(payload.get("reason") or "user rejected tool call")
    approval.resolved_by = current_user.id
    approval.resolved_at = now
    tool = await db.scalar(
        select(ResponseToolExecution).where(
            ResponseToolExecution.response_id == response_id,
            ResponseToolExecution.call_id == approval.call_id,
        )
    )
    if tool:
        tool.status = approval.status
        tool.error_message = approval.reason
    resolved_event = await append_event(
        db,
        response_id=response_id,
        event_type="opentrace.approval.resolved",
        payload={
            "approval_id": approval.id,
            "call_id": approval.call_id,
            "approved": approved,
            "status": approval.status,
            "operation_class": getattr(
                approval,
                "operation_class",
                getattr(approval, "side_effect_level", "write"),
            ),
            "required_approvals": required_approvals,
            "received_approvals": received_approvals,
        },
    )
    # 批准和拒绝都恢复 Manager；拒绝会作为类型化工具结果进入后续解释。
    _prepare_response_for_approval_resume(response)
    await _ensure_approval_resume_outbox(
        db,
        response_id=response_id,
        approval_id=approval.id,
    )
    await db.commit()
    outcome = "approved" if approved else "rejected"
    RESPONSE_APPROVAL_DECISIONS_TOTAL.labels(
        outcome=outcome,
        required_approvals=str(required_approvals),
    ).inc()
    created_at = getattr(approval, "created_at", None)
    if created_at is not None:
        normalized_created_at = (
            created_at.replace(tzinfo=UTC)
            if created_at.tzinfo is None
            else created_at.astimezone(UTC)
        )
        RESPONSE_APPROVAL_RESOLUTION_DURATION.labels(
            outcome=outcome,
            required_approvals=str(required_approvals),
        ).observe(max(0.0, (now - normalized_created_at).total_seconds()))
    return {
        "approved": approved,
        "status": approval.status,
        "call_id": approval.call_id,
        "approval_id": approval.id,
        "required_approvals": required_approvals,
        "received_approvals": received_approvals,
        "current_user_decision": "approved" if approved else "rejected",
        "starting_after": resolved_event.sequence_number,
    }


@router.post("/responses/{response_id}/tool-approvals")
async def approve_response_tool(
    response_id: str,
    request: Request,
    payload: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Approve/reject a side-effecting tool call and append an auditable event.

    The endpoint is intentionally idempotent: repeating the same decision
    returns the existing ledger row and never executes a tool twice.
    """
    call_id = str(payload.get("call_id") or "")
    if not call_id:
        raise AppException(ErrorCodes.PARAM_INVALID.code, message="缺少工具调用标识")
    return await _resolve_response_tool_approval(
        response_id=response_id,
        request=request,
        payload=payload,
        current_user=current_user,
        db=db,
        call_id=call_id,
    )


@router.post("/responses/{response_id}/approvals/{approval_id}/resolve")
async def resolve_response_approval(
    response_id: str,
    approval_id: str,
    request: Request,
    payload: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    return await _resolve_response_tool_approval(
        response_id=response_id,
        request=request,
        payload=payload,
        current_user=current_user,
        db=db,
        approval_id=approval_id,
    )


@router.get("/response-approvals/pending")
async def list_pending_four_eye_approvals(
    request: Request,
    limit: int = 100,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> dict:
    """列出同 Scope 的高风险待复核项，不暴露 Response 正文或原始敏感参数。"""

    normalized_role = str(current_user.role or "").strip().lower()
    if not (is_enterprise_admin(current_user) or normalized_role == "sre"):
        raise AppException(ErrorCodes.PERMISSION_DENIED.code, message="当前角色无权复核生产变更")
    tenant_id, workspace_id = _scope(request, current_user)
    bounded_limit = max(1, min(int(limit), 200))
    rows = await db.execute(
        select(ResponseApproval, ResponseRecord.user_id, ResponseRecord.created_at)
        .join(ResponseRecord, ResponseRecord.id == ResponseApproval.response_id)
        .where(
            ResponseRecord.tenant_id == tenant_id,
            ResponseRecord.workspace_id == workspace_id,
            ResponseRecord.status == "requires_action",
            ResponseApproval.status.in_(("pending", "pending_secondary")),
            ResponseApproval.required_approvals >= 2,
        )
        .order_by(ResponseApproval.created_at)
        .limit(bounded_limit)
    )
    items = []
    for approval, requested_by, response_created_at in rows.all():
        decisions = [
            dict(item)
            for item in approval.approval_decisions or []
            if isinstance(item, dict) and str(item.get("user_id") or "")
        ]
        current_decision = next(
            (
                "approved" if item.get("approved") is True else "rejected"
                for item in decisions
                if str(item.get("user_id")) == current_user.id
            ),
            None,
        )
        items.append(
            {
                "id": approval.id,
                "response_id": approval.response_id,
                "call_id": approval.call_id,
                "tool_name": approval.tool_name,
                "side_effect": approval.side_effect_level,
                "operation_class": approval.operation_class,
                "arguments": mask_sensitive(dict(approval.arguments or {})),
                "status": approval.status,
                "required_approvals": int(approval.required_approvals or 1),
                "received_approvals": len(
                    {str(item["user_id"]) for item in decisions if item.get("approved") is True}
                ),
                "current_user_decision": current_decision,
                "requested_by": requested_by,
                "response_created_at": (
                    response_created_at.isoformat() if response_created_at else None
                ),
                "created_at": approval.created_at.isoformat() if approval.created_at else None,
            }
        )
    return {"items": items, "limit": bounded_limit}


@router.get("/messages/{message_id}/versions")
async def message_versions(
    message_id: str,
    request: Request,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    item = await db.scalar(select(ResponseItem).where(ResponseItem.id == message_id))
    response_id = item.response_id if item else message_id
    response = await _owned_response(
        db, response_id=response_id, user=current_user, request=request
    )
    if response is None:
        return []
    siblings = (
        (
            await db.execute(
                select(ResponseRecord)
                .where(
                    ResponseRecord.conversation_id == response.conversation_id,
                    ResponseRecord.parent_response_id == response.parent_response_id,
                    ResponseRecord.user_id == current_user.id,
                    ResponseRecord.tenant_id == response.tenant_id,
                    ResponseRecord.workspace_id == response.workspace_id,
                )
                .order_by(ResponseRecord.created_at)
            )
        )
        .scalars()
        .all()
    )
    return [
        {
            "id": row.id,
            "response_id": row.id,
            "status": row.status,
            "model": row.model,
            "active": row.id == response.id,
            "created_at": row.created_at.isoformat(),
        }
        for row in siblings
    ]


@router.post("/responses/{response_id}/feedback")
async def response_feedback(
    response_id: str,
    request: Request,
    payload: dict,
    current_user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    item = await db.scalar(select(ResponseItem).where(ResponseItem.id == response_id))
    canonical_id = item.response_id if item else response_id
    response = await _owned_response(
        db, response_id=canonical_id, user=current_user, request=request
    )
    if response is None:
        return {"status": "not_found"}
    response_items = (
        (
            await db.execute(
                select(ResponseItem)
                .where(ResponseItem.response_id == response.id)
                .order_by(ResponseItem.sequence_number)
            )
        )
        .scalars()
        .all()
    )
    query = next(
        (row.content or "" for row in response_items if row.item_type == "input_message"), ""
    )
    answer = next(
        (row.content or "" for row in reversed(response_items) if row.item_type == "message"), ""
    )
    db.add(
        Feedback(
            id=str(uuid.uuid4()),
            session_id=response.conversation_id,
            query=query,
            response=answer,
            feedback_type=str(payload.get("feedback_type") or "none"),
            score=payload.get("score"),
            correction=payload.get("correction"),
            feedback_metadata=json.dumps(
                {"source": "responses_api", "response_id": response.id, "user_id": current_user.id},
                ensure_ascii=False,
            ),
        )
    )
    data_feedback_id = None
    assistant_item = next(
        (row for row in reversed(response_items) if row.item_type == "message"), None
    )
    assistant_payload = dict(assistant_item.payload or {}) if assistant_item else {}
    data_answer = dict(assistant_payload.get("data_answer") or {})
    run_id = str(data_answer.get("data_agent_run_id") or "").strip()
    if run_id:
        run_record = await db.scalar(
            select(DataAgentRunRecord).where(
                DataAgentRunRecord.id == run_id,
                DataAgentRunRecord.user_id == current_user.id,
                DataAgentRunRecord.tenant_id == response.tenant_id,
                DataAgentRunRecord.workspace_id == response.workspace_id,
            )
        )
        if run_record is not None:
            feedback_type = str(payload.get("feedback_type") or "none").lower()
            score = payload.get("score")
            verdict = None
            if feedback_type in {"like", "helpful", "correct"}:
                verdict = "correct"
            elif feedback_type in {"dislike", "unhelpful", "incorrect"}:
                verdict = "incorrect"
            elif isinstance(score, int | float):
                verdict = "correct" if score >= 0.5 else "incorrect"
            if verdict is None:
                await db.commit()
                return {
                    "status": "accepted",
                    "response_id": response.id,
                    "data_agent_feedback_id": None,
                }
            data_feedback_id = str(uuid.uuid4())
            positive_feedback = verdict == "correct"
            resolved_at = datetime.now(UTC) if positive_feedback else None
            db.add(
                DataAgentFeedback(
                    id=data_feedback_id,
                    run_id=run_id,
                    user_id=current_user.id,
                    tenant_id=response.tenant_id,
                    workspace_id=response.workspace_id,
                    verdict=verdict,
                    candidate_id=run_record.selected_candidate_id,
                    corrected_sql=str(payload.get("corrected_sql") or "") or None,
                    comment=str(payload.get("correction") or "") or None,
                    metadata_json={
                        "source": "responses_feedback_bridge",
                        "response_id": response.id,
                        "response_item_id": assistant_item.id if assistant_item else None,
                        "feedback_type": feedback_type,
                        "score": score,
                    },
                    status="resolved" if positive_feedback else "open",
                    resolution_note="正向反馈自动归档" if positive_feedback else None,
                    resolved_by=current_user.id if positive_feedback else None,
                    resolved_at=resolved_at,
                )
            )
            if settings.data_agent_learning_enabled:
                scope = DataScope(
                    user_id=current_user.id,
                    tenant_id=response.tenant_id,
                    workspace_id=response.workspace_id,
                    data_source_id=run_record.data_source_id,
                )
                run = await OpenTraceRunRepository(db).get(run_id, scope)
                if run is not None:
                    await OpenTraceLearningRepository(db).record_feedback(
                        run,
                        verdict=verdict,
                        candidate_id=run_record.selected_candidate_id,
                        corrected_sql=str(payload.get("corrected_sql") or "") or None,
                    )
    await db.commit()
    return {
        "status": "accepted",
        "response_id": response.id,
        "data_agent_feedback_id": data_feedback_id,
    }
